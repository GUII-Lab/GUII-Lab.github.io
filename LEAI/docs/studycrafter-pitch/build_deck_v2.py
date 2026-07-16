"""Build the LEAI pitch deck (v2, post-review) as a .pptx.

Writes studycrafter-pitch-v2.pptx next to this file. The original
studycrafter-pitch.pptx is left untouched.

Run:
    cd LEAI/docs/studycrafter-pitch
    uv run --with python-pptx --with pillow python build_deck_v2.py

Art sources
-----------
NEW captures (v0.2.8, cropped, no browser chrome, no sidebar), taken from the
dev stack against `test-course`:
    SHOTS/s3-setup-v2.png       survey setup: templates + bot persona
    SHOTS/s6-claims-trio.png    three Insights claims w/ student-count chips

INTERIM art, cropped out of the previous deck's media (student-facing pages
have no sidebar and no version footer, so they are not visibly stale; they do
NOT yet show the revise/add-to hint):
    SHOTS/old-image5.png        consent modal  -> slide 4
    SHOTS/old-image4.png        form chat, "Area 2 of 4" -> slide 5
    SHOTS/old-image10.png       feedback chat  -> backup slide
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "shots"
OUT = HERE / "studycrafter-pitch-v2.pptx"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NAVY = RGBColor(0x0F, 0x23, 0x4E)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0xC2, 0x41, 0x0C)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_FG = RGBColor(0x92, 0x40, 0x0E)
LINE = RGBColor(0xE2, 0xE8, 0xF0)

EMU_PER_IN = 914400


def inches(v):
    return Emu(int(v * EMU_PER_IN))


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, italic=False, line=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if line:
        p.line_spacing = line
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rich(slide, left, top, width, height, chunks, *, size=16, line=1.35):
    """chunks: list of list-of-(text, bold, color) -> one paragraph per row."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, row in enumerate(chunks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line
        p.space_after = Pt(7)
        for (txt, bold, color) in row:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def add_band(slide, left, top, width, height, fill=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_image_fitted(slide, path, left, top, max_w, max_h, *, border=True):
    img = Image.open(path)
    iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    pic = slide.shapes.add_picture(str(path), cx, cy, width=w, height=h)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- slides

def slide_1(prs):
    s = blank(prs)
    add_text(s, inches(0.9), inches(1.9), inches(11.5), inches(1.0),
             "LEAI — Learning Experience AI",
             size=54, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, inches(1.35), inches(3.15), inches(10.6), inches(0.55),
             "Course feedback as a 5-minute chat, not a form.",
             size=24, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_band(s, inches(2.2), inches(4.75), inches(8.9), inches(0.95))
    add_text(s, inches(2.2), inches(5.03), inches(8.9), inches(0.5),
             "Instructor sets up   →   Student chats   →   Instructor reads",
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, inches(0), inches(6.95), inches(13.33), inches(0.3),
             "GUII Lab · UC Santa Cruz", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    notes(s, """Hi everyone — I want to show you a tool from our lab called LEAI, Learning Experience AI.

The one-line version is on the slide: course feedback as a five-minute chat instead of a form.

Slightly longer: students open a link and have a short conversation with an AI — it asks follow-up questions the way a TA would. The instructor gets back a summary they can skim in about two minutes, and every claim in it traces back to the student who said it.

[If StudyHelper comes up] Different tool, different job: StudyHelper teaches during the course. LEAI asks about the course. LEAI is not a StudyHelper add-on and doesn't depend on it.

Three parts, and I'll show you each: the instructor sets up a survey, the student chats, the instructor reads what came out.""")


def slide_2(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.7),
             "Why LEAI — the gap in course feedback", size=36, bold=True, color=NAVY)
    top, h, w, gap = inches(1.75), inches(2.5), inches(3.8), inches(0.22)
    lefts = [inches(0.65), inches(0.65) + w + gap, inches(0.65) + 2 * (w + gap)]
    cols = [
        ("End-of-term evals", "Too late"),
        ("Mid-semester Likert", "Too shallow"),
        ("Open-ended writing", "Too much"),
    ]
    for left, (lbl, hd) in zip(lefts, cols):
        add_band(s, left, top, w, h)
        add_text(s, left + inches(0.25), top + inches(0.65), w - inches(0.5), inches(0.4),
                 lbl, size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, left + inches(0.25), top + inches(1.25), w - inches(0.5), inches(0.6),
                 hd, size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    pt = inches(4.85)
    add_band(s, inches(0.65), pt, inches(12), inches(1.2), fill=NAVY)
    add_text(s, inches(0.95), pt + inches(0.35), inches(11.4), inches(0.5),
             "LEAI sits in the middle", size=26, bold=True, color=WHITE)
    notes(s, """Existing tools force a tradeoff between timing, depth, and workload. Three ways to collect course feedback today, and each one fails differently.

End-of-term evals — too late. By the time results land, the students who took the class are already gone. Nothing you learn helps the people who told you.

Mid-semester Likert — too shallow. Fast to fill out, but a 3-out-of-5 tells you what, never why. You can't act on it.

Open-ended writing — too much. It's rich, but at forty long replies a week nobody has time to read them all and turn them into action. Students know that, so they stop trying.

LEAI is our attempt at the middle path: five minutes for the student, two minutes for the instructor, and every claim traceable back to the student who said it.

[This slide is load-bearing. Do not cut it for time. The results slide has no baseline without it.]""")


def slide_3(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.7),
             "Set up a survey in minutes", size=36, bold=True, color=NAVY)
    add_text(s, inches(0.65), inches(1.2), inches(12), inches(0.5),
             "Pick a template.   Set how the AI talks.   Share a link.",
             size=19, color=MUTED, italic=True)

    add_image_fitted(s, SHOTS / "s3-setup-v2.png",
                     inches(0.65), inches(1.95), inches(6.5), inches(4.5))

    rx = inches(7.6)
    add_text(s, rx, inches(2.6), inches(5.1), inches(0.4),
             "Three modes", size=20, bold=True, color=ACCENT)
    add_rich(s, rx, inches(3.25), inches(5.1), inches(2.0), [
        [("Free-Form", True, NAVY)],
        [("In-Group", True, NAVY)],
        [("Structured Reflection", True, NAVY)],
    ], size=20, line=1.4)
    notes(s, """This is the instructor side — PromptDesigner.

You pick a template — Weekly, Midterm, Project, Exit Ticket — so you're not starting from a blank page. You give the bot a name and a role; that's what shapes how it talks to students, and the name is what they see at the top of their chat. Then you share a link. That's the whole setup — minutes, not hours. One per week is a common rhythm.

Three modes, and gloss each one as you say it — don't assume the names explain themselves:
- Free-Form — open feedback. The default.
- In-Group — students rate and reflect on their team. (Say "rate your team" out loud — "In-Group" collides with the social-psych in-group/out-group sense and people will hear the wrong thing.)
- Structured Reflection — the AI walks every section in order. That's the next slide but one.

[Screenshot is seeded demo data on a test course.]""")


def slide_4(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.7),
             "What students experience", size=36, bold=True, color=NAVY)
    add_text(s, inches(0.65), inches(1.2), inches(12), inches(0.5),
             "A five-minute chat. No login, no app — just a link.",
             size=19, color=MUTED, italic=True)

    add_image_fitted(s, SHOTS / "old-image5.png",
                     inches(0.65), inches(1.95), inches(6.5), inches(4.3))

    rx = inches(7.6)
    add_rich(s, rx, inches(2.8), inches(5.1), inches(2.4), [
        [("Conversational", True, NAVY)],
        [("Anonymous", True, NAVY), ("  (per survey)", False, MUTED)],
        [("Voice or typing", True, NAVY)],
    ], size=20, line=1.4)
    notes(s, """This is what a student sees. They open a link — no login, no app, no account.

First screen is a quick consent: they agree to the terms, and there's a separate optional box to let our lab use the anonymised data for research. That one is optional and it's per session — they can change it every time.

Then it's just a chat. The AI asks a question, they answer, and it follows up — that's the part a form can't do. They can talk instead of type, which matters when they're tired or on a phone.

On anonymity: it's the instructor's choice, set per survey — anonymous, or with a code. I say it that way because it IS a setting; don't claim students are always anonymous.

LEAI doesn't grade. It collects.""")


def slide_5(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.75),
             "Structured Reflection — when every section must be answered",
             size=30, bold=True, color=NAVY)

    add_image_fitted(s, SHOTS / "old-image4.png",
                     inches(0.65), inches(1.5), inches(6.5), inches(4.6))

    rx = inches(7.6)
    add_rich(s, rx, inches(2.5), inches(5.1), inches(2.6), [
        [("The AI walks every section in order.", True, NAVY)],
        [("Students can revise any answer, anytime.", True, NAVY)],
    ], size=20, line=1.45)
    notes(s, """Structured Reflection. Same chat, but guided.

On the left is the student. It still looks like a normal chat — but the header says "Area 2 of 4," so the AI is walking them through every section in order. They finish when every section is covered.

Why it exists: sometimes you actually need every section answered — a graded reflection where every part of the rubric matters. So the structure lives in the engine, not on the screen.

One thing we added after student feedback: students can revise or add to any earlier answer at any time — they just say "actually, change what I said about..." Students told us the old flow felt like being marched through a form. This is the fix.

[Honesty: instructors can't author these outlines in-app yet — we set them up. That's why I say "picks from a set of outlines we've set up," and it's on the ideas slide.]

[Interim art: this capture predates the revise hint, so the hint isn't visible in the screenshot yet.]""")


def slide_6(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(11), inches(0.7),
             "From conversations to insights", size=36, bold=True, color=NAVY)

    # demo-data stamp
    st = add_band(s, inches(11.0), inches(0.5), inches(1.65), inches(0.36), fill=AMBER_BG)
    st.line.color.rgb = RGBColor(0xFC, 0xD3, 0x4D)
    st.line.width = Pt(0.75)
    add_text(s, inches(11.0), inches(0.585), inches(1.65), inches(0.25),
             "DEMO DATA", size=10, bold=True, color=AMBER_FG, align=PP_ALIGN.CENTER)

    add_text(s, inches(0.65), inches(1.18), inches(11), inches(0.45),
             "Every claim shows how many students back it.",
             size=21, bold=True, color=ACCENT)

    # Fills the width: 2-claim crop (aspect ~2.44) at 11.9in wide ≈ 4.87in tall.
    add_image_fitted(s, SHOTS / "s6-claims-pair.png",
                     inches(0.65), inches(1.95), inches(11.9), inches(4.6))
    notes(s, """Now the instructor side — the Feedback Analyzer.

Free-form chat is rich but messy. Forty conversations is a lot to read, and for a big class it's hopeless. So we give you three views and let you triangulate — three views, so you don't have to trust one summary:
- Instructor Insights — the AI's read, on screen now
- Keyness — the words that stand out compared to a typical class. Gloss it; nobody outside corpus linguistics knows the word. This one uses no AI at all, it's statistics.
- Raw responses — one click to the actual student text

And the other half of the change: claims only one student made get quarantined into their own "Mentioned by one student" group, instead of being laundered into "some students feel..."

Look at these three bullets. Five students are overwhelmed by the stats side. Twelve found the framing clear and valuable. Seven found it too basic and want to go faster. That's not one summary telling you "students had mixed feelings" — that's a real disagreement, with sizes, and you can see who's on each side.

That's the change we made this quarter. Every claim carries a student count, and claims only one student made get quarantined into their own group instead of being laundered into "some students feel..."

[If asked — are those counts generated by the model?] No. The count is computed: it's the number of distinct responses the model cited, deduped, and any citation that doesn't resolve to a real response is dropped, so a hallucinated ID can't inflate it. The honest caveat is that code does the counting but the model chooses what to cite — so a count can be exact while a citation is wrong. Every one is a click from the raw text, which is how you'd catch it.

[Numbers on this screenshot are seeded demo data on a test course — the real numbers are on the next slide. Say this out loud once.]

[If Feedback Chat comes up] There's also a chat interface where you can interrogate the whole set of responses and it cites what it uses — happy to show it after.""")


def slide_7(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.4), inches(12), inches(0.6),
             "What happened when we deployed it", size=32, bold=True, color=NAVY)

    add_text(s, inches(0.65), inches(1.5), inches(12), inches(1.3),
             "22 of 25  →  25 of 25", size=72, bold=True, color=NAVY)
    add_text(s, inches(0.65), inches(2.85), inches(12), inches(0.4),
             "Week 9 → Week 10, HCI 271", size=20, color=MUTED)

    add_band(s, inches(0.65), inches(3.6), inches(12), inches(0.95))
    add_text(s, inches(0.9), inches(3.88), inches(11.5), inches(0.45),
             "Same assignment, same credit — students freely picked LEAI or the "
             "instructor's PDF form.", size=18, color=INK)

    add_rich(s, inches(0.65), inches(5.0), inches(12), inches(0.9), [
        [("She used Insights to find students confused about co-design methods — "
          "and added new lecture content that quarter.", True, INK)],
    ], size=18, line=1.3)

    add_text(s, inches(0.65), inches(6.35), inches(12), inches(0.35),
             "3 courses · 3 instructors    |    Usability (SUS): 65.0 / 67.5, n=6 — unpowered",
             size=13, color=MUTED)
    notes(s, """This is the part worth your time.

Three courses this quarter — HCI 271 with Magy, CMPM 230 with Reza, HCI 220 with Mahnaz. Two grad HCI, one games. All three at UCSC, same department, overlapping instructor pool — that's one context sampled three times, not three independent contexts.

The number on the left is the one I'd point at. In HCI 271, weeks 9 and 10, students could choose: use LEAI, or submit through the PDF form the instructor originally designed. Same assignment, same credit, free pick. Week 9, 22 of 25 chose LEAI. Week 10, all 25.

Be precise about what that does and doesn't show. It's the same 25 students both weeks, so 22-to-25 is three people — that's not a trend, don't sell it as one. And "chose it" is not "preferred it": a link is less friction than a download-fill-submit PDF, so some of this is effort. What it does show is that essentially nobody opted out when opting out was free — it cleared the acceptability bar in a real graded course.

The line on the slide is the one I'd actually hang the case on, and it's behavioural: Magy used Insights to find students were confused about co-design methods, and added new lecture content that same quarter. Feedback becoming a teaching change while the course was still running — that's the thing end-of-term evals structurally cannot do. Behaviour beats testimony.

Magy's other read, if it's useful: reflections through LEAI were "a lot more authentic" than the form, where answers were often too similar — she thinks copying. That's her hypothesis, not a measurement; we haven't run the text-similarity comparison that would test it, and we could. Worth saying plainly: Magy is the instructor here AND our lab director, so this is an instructor's report, not an independent evaluation. Say that before someone else thinks it.

Reza (CMPM 230): said LEAI was better than the current system, and would use again. "Said" — not "rated." No instrument was administered; don't dress a conversation up as a scale.

Mahnaz (HCI 220): small number of students, and we never got a formal interview — feedback came from ad-hoc Zoom conversations. I'm counting the course in "three," so I should report the thinness rather than let it disappear.

On SUS: 65.0 and 67.5, both under the 68 benchmark, n=6 out of ~25 and self-selected — about a 24% response rate. Unpowered; I won't over-read it. But the interesting part is the tension: adoption was near-total while usability scored average-to-below. So adoption here isn't tracking usability — it's an acceptability signal, maybe an effort signal. That's a finding, not an embarrassment.""")


def slide_8(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.7),
             "Ideas we're weighing — not a roadmap", size=32, bold=True, color=NAVY)

    t1 = inches(1.6)
    add_band(s, inches(0.65), t1, inches(12), inches(2.25))
    add_text(s, inches(0.9), t1 + inches(0.22), inches(11.5), inches(0.3),
             "Exploring", size=14, bold=True, color=ACCENT)
    add_rich(s, inches(0.9), t1 + inches(0.7), inches(11.3), inches(1.1), [
        [("Student control over the conversation flow", True, INK)],
        [("Instructors authoring their own outlines", True, INK)],
        [("Prompt sharing between instructors", True, INK)],
    ], size=18, line=1.35)

    t2 = inches(4.1)
    b2 = add_band(s, inches(0.65), t2, inches(12), inches(1.0), fill=WHITE)
    b2.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    b2.line.width = Pt(1)
    add_text(s, inches(0.9), t2 + inches(0.18), inches(11.5), inches(0.3),
             "Asked for, not planned", size=14, bold=True, color=MUTED)
    add_text(s, inches(0.9), t2 + inches(0.55), inches(11.3), inches(0.35),
             "Canvas integration", size=18, bold=True, color=INK)

    t3 = inches(5.35)
    add_band(s, inches(0.65), t3, inches(12), inches(1.05), fill=NAVY)
    add_text(s, inches(0.9), t3 + inches(0.3), inches(11.5), inches(0.5),
             "Want to try it in your course next quarter? Come find me.",
             size=22, bold=True, color=WHITE)
    notes(s, """Where this could go. These are ideas we're weighing — not commitments, and not a roadmap. I'm saying that because a list like this gets remembered as a promise.

The one with the most evidence behind it is student control over the conversation flow. An instructor asked for it, and students asked for it independently in their survey write-ins. When two sources that never talked to each other say the same thing, that's worth listening to.

Instructors authoring their own outlines — right now we set those up by hand, which doesn't scale past a few courses. An LLM-assisted survey wizard is the obvious companion.

Prompt sharing — instructors keep reinventing the same weekly check-in.

Canvas: it's the single most common request, and it is not started. I'm listing it so nobody thinks I haven't heard it, not to signal it's coming.

And the ask: if you want to try this in your course next quarter, come find me. Setup is a few minutes a week; it's five minutes of your students' time.""")


def slide_backup(prs):
    s = blank(prs)
    add_text(s, inches(0.65), inches(0.45), inches(12), inches(0.7),
             "Feedback Chat — ask questions of the responses",
             size=30, bold=True, color=NAVY)
    add_text(s, inches(0.65), inches(1.2), inches(12), inches(0.4),
             "Backup slide — shown only if asked.", size=15, color=MUTED, italic=True)
    add_image_fitted(s, SHOTS / "old-image10.png",
                     inches(1.4), inches(1.85), inches(10.5), inches(4.5))
    notes(s, """Backup. Only if someone asks.

Instructors told me they wanted to just talk to the data instead of reading it. So: ask a question — "what are students struggling with this week?" — and it answers with numbered citations back to the specific responses it used. You can click any one and read the student's actual words.

You can also open the system prompt and edit it. We thought that mattered for trust: you should be able to see what's under the hood.

[This is a way of using the Analyzer, not a fourth product — don't let it read as a separate tool.]""")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_backup(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
