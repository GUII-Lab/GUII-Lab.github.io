"""Build the StudyCrafter LEAI pitch deck as a .pptx file.

Run:  uv run --with python-pptx python build_deck.py

Output: studycrafter-pitch.pptx in this directory. Drag it into Google
Drive and double-click to open in Google Slides.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

REPO = Path(__file__).resolve().parents[3]
ASSETS = REPO / "LEAI" / "guide-assets"
STUDENT_ASSETS = REPO / "LEAI" / "student-guide-assets"
OUT = Path(__file__).resolve().parent / "studycrafter-pitch.pptx"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NAVY = RGBColor(0x0F, 0x23, 0x4E)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0xC2, 0x41, 0x0C)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_band(slide, left, top, width, height, fill=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_image_fitted(slide, path, left, top, max_w, max_h):
    """Insert an image scaled to fit the box while preserving aspect ratio."""
    from PIL import Image

    img = Image.open(path)
    iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(str(path), cx, cy, width=w, height=h)


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def fill_background(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back so subsequent shapes draw over it.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def slide_1(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(900000), Emu(1600000), Emu(10400000), Emu(900000),
             "LEAI — Learning Experience AI",
             size=54, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s, Emu(1400000), Emu(2700000), Emu(9400000), Emu(1100000),
             "A sibling lab tool to StudyHelper.\n"
             "StudyHelper helps students learn. LEAI listens to how it went.",
             size=22, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    add_band(s, Emu(900000), Emu(4500000), Emu(10400000), Emu(1200000), fill=SOFT)
    add_text(s, Emu(900000), Emu(4800000), Emu(10400000), Emu(700000),
             "PromptDesigner   →   FeedbackCollector   →   FeedbackAnalyzer",
             size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s, Emu(0), Emu(6450000), SLIDE_W, Emu(300000),
             "GUII Lab · UC Santa Cruz",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_motivation(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(600000), Emu(400000), Emu(11000000), Emu(700000),
             "Why LEAI — the gap in course feedback",
             size=36, bold=True, color=NAVY)
    add_text(s, Emu(600000), Emu(1100000), Emu(11000000), Emu(500000),
             "Existing tools force a tradeoff between timing, depth, and workload.",
             size=20, color=MUTED, italic=True)

    col_top = Emu(1900000)
    col_h = Emu(2700000)
    col_w = Emu(3500000)
    gap = Emu(200000)
    col_left_1 = Emu(600000)
    col_left_2 = col_left_1 + col_w + gap
    col_left_3 = col_left_2 + col_w + gap

    def problem_col(left, label, headline, body):
        add_band(s, left, col_top, col_w, col_h, fill=SOFT)
        add_text(s, left + Emu(250000), col_top + Emu(220000),
                 col_w - Emu(500000), Emu(400000),
                 label, size=14, bold=True, color=ACCENT)
        add_text(s, left + Emu(250000), col_top + Emu(680000),
                 col_w - Emu(500000), Emu(600000),
                 headline, size=18, bold=True, color=NAVY)
        add_text(s, left + Emu(250000), col_top + Emu(1400000),
                 col_w - Emu(500000), Emu(1200000),
                 body, size=14, color=INK)

    problem_col(col_left_1, "End-of-term evals",
                "Too late to help",
                "By the time results land, the students who took the class are already gone.")
    problem_col(col_left_2, "Mid-semester Likert",
                "Fast but shallow",
                "Quick to fill out, but the signal is thin — you cannot tell what is actually going wrong.")
    problem_col(col_left_3, "Open-ended writing",
                "Rich but unreadable",
                "Forty long replies a week. No one has time to read them all and turn them into action.")

    panel_top = Emu(4900000)
    panel_h = Emu(1500000)
    add_band(s, Emu(600000), panel_top, Emu(11000000), panel_h, fill=NAVY)
    add_text(s, Emu(800000), panel_top + Emu(180000), Emu(10600000), Emu(500000),
             "LEAI — the middle path",
             size=22, bold=True, color=WHITE)
    add_text(s, Emu(800000), panel_top + Emu(700000), Emu(10600000), Emu(700000),
             "5 minutes for the student. 2 minutes for the instructor. Every claim "
             "traceable back to the student who said it.",
             size=16, color=WHITE)

    add_text(s, Emu(600000), Emu(6500000), Emu(11000000), Emu(280000),
             "StudyHelper helps students learn. LEAI listens to how it went.",
             size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


def slide_2(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(600000), Emu(400000), Emu(11000000), Emu(700000),
             "Set up a survey in minutes",
             size=36, bold=True, color=NAVY)
    add_text(s, Emu(600000), Emu(1100000), Emu(11000000), Emu(500000),
             "Pick a template.   Write the bot's voice.   Share a link.",
             size=20, color=MUTED, italic=True)

    img_top = Emu(1900000)
    img_h = Emu(3800000)
    img_w = Emu(5400000)
    add_image_fitted(s, ASSETS / "test-f1-f2-survey-list.png",
                     Emu(400000), img_top, img_w, img_h)
    add_image_fitted(s, ASSETS / "test-f11-edit-modal.png",
                     Emu(6400000), img_top, img_w, img_h)

    add_text(s, Emu(600000), Emu(6000000), Emu(11000000), Emu(500000),
             "Three modes:  Free-Form  ·  In-Group  ·  Structured Reflection",
             size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_3(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(600000), Emu(400000), Emu(11000000), Emu(700000),
             "What students experience",
             size=36, bold=True, color=NAVY)

    add_text(s, Emu(600000), Emu(1150000), Emu(11000000), Emu(450000),
             "Conversational  ·  Anonymous  ·  Group-aware  ·  Voice-input",
             size=20, bold=True, color=ACCENT)

    img_top = Emu(1800000)
    img_h = Emu(3300000)
    img_w = Emu(5400000)
    add_image_fitted(s, ASSETS / "feedback-consent-modal.png",
                     Emu(400000), img_top, img_w, img_h)
    add_image_fitted(s, ASSETS / "chat-new-session.png",
                     Emu(6400000), img_top, img_w, img_h)

    thumb_top = Emu(5300000)
    thumb_h = Emu(1300000)
    thumb_w = Emu(3500000)
    add_image_fitted(s, ASSETS / "feedback-header-anonymity.png",
                     Emu(400000), thumb_top, thumb_w, thumb_h)
    add_image_fitted(s, STUDENT_ASSETS / "group-team-select.png",
                     Emu(4350000), thumb_top, thumb_w, thumb_h)
    add_image_fitted(s, STUDENT_ASSETS / "mic-recording.png",
                     Emu(8300000), thumb_top, thumb_w, thumb_h)


def slide_4(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(600000), Emu(400000), Emu(11000000), Emu(700000),
             "Structured Reflection — when every section must be answered",
             size=30, bold=True, color=NAVY)

    img_top = Emu(1300000)
    img_h = Emu(3600000)
    img_w = Emu(5400000)
    add_image_fitted(s, ASSETS / "pd-form-mode.png",
                     Emu(400000), img_top, img_w, img_h)
    add_image_fitted(s, ASSETS / "feedback-form-mode.png",
                     Emu(6400000), img_top, img_w, img_h)

    panel_top = Emu(5100000)
    panel_h = Emu(1500000)
    add_band(s, Emu(600000), panel_top, Emu(11000000), panel_h, fill=SOFT)

    box = add_text(s, Emu(800000), panel_top + Emu(150000), Emu(10600000), Emu(1200000),
                   "What:  ", size=16, bold=True, color=ACCENT)
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AI walks every section in order."
    r.font.size = Pt(16)
    r.font.color.rgb = INK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r1 = p2.add_run()
    r1.text = "Why:  "
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p2.add_run()
    r2.text = "When you need every section answered — for example, a graded team reflection."
    r2.font.size = Pt(16)
    r2.font.color.rgb = INK

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    r1 = p3.add_run()
    r1.text = "How:  "
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p3.add_run()
    r2.text = "Instructor picks a schema in Prompt Designer. Student sees a normal chat."
    r2.font.size = Pt(16)
    r2.font.color.rgb = INK


def slide_5(prs):
    s = blank_slide(prs)
    fill_background(s, WHITE)

    add_text(s, Emu(600000), Emu(400000), Emu(11000000), Emu(700000),
             "From conversations to insights",
             size=36, bold=True, color=NAVY)
    add_text(s, Emu(600000), Emu(1100000), Emu(11000000), Emu(500000),
             "Quicktake  ·  Keyness  ·  Raw responses",
             size=20, bold=True, color=ACCENT)

    img_top = Emu(1700000)
    img_h = Emu(3000000)
    img_w = Emu(5400000)
    add_image_fitted(s, ASSETS / "analyzer-quicktake.png",
                     Emu(400000), img_top, img_w, img_h)
    add_image_fitted(s, ASSETS / "analyzer-keyness.png",
                     Emu(6400000), img_top, img_w, img_h)

    thumb_w = Emu(2000000)
    thumb_h = Emu(900000)
    add_image_fitted(s, ASSETS / "analyzer-ngram-clickthrough.png",
                     Emu(5100000), Emu(4750000), thumb_w, thumb_h)

    panel_top = Emu(5800000)
    panel_h = Emu(800000)
    add_band(s, Emu(600000), panel_top, Emu(11000000), panel_h, fill=SOFT)
    add_text(
        s, Emu(800000), panel_top + Emu(150000), Emu(10600000), Emu(550000),
        "Free-form chat is rich but messy. We did not want instructors to trust one black-box "
        "summary. Three views, so you can triangulate.",
        size=14, color=INK, italic=True,
    )

    add_text(s, Emu(0), Emu(6620000), SLIDE_W - Emu(400000), Emu(220000),
             "More at LEAI/instructor-guide.html",
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1(prs)
    slide_motivation(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
